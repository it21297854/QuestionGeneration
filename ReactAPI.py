import fitz
import spacy
import random
import os
from flask import Flask, request, jsonify
from flask_cors import CORS
from docx import Document
from pptx import Presentation

app = Flask(__name__)
CORS(app)

nlp = spacy.load("en_core_web_sm")

UPLOAD_FOLDER = "uploads"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)

# Store questions
questions_store = []

# IT-related words for each level
high_level_keywords = [
    'machine_learning', 'natural language processing', 'optical character recognition',
    'deep learning', 'cybersecurity', 'artificial intelligence', 'quantum computing',
    'blockchain', 'predictive models', 'threat detection'
]

medium_level_keywords = [
    'flask', 'mysql', 'sqlalchemy', 'bootstrap', 'figma', 'role-based access control',
    'data management', 'database schema', 'security testing', 'performance optimization',
    'horizontal scaling', 'devops', 'docker', 'cloud computing', 'encryption',
    'referential integrity', 'modular architecture'
]

low_level_keywords = [
    'modular design', 'user interface', 'responsive design', 'mobile compatibility',
    'authentication', 'authorization', 'secure authentication', 'system performance',
    'UI/UX principles', 'data consistency', 'software architecture', 'scalability',
    'secure access control', 'index numbers', 'real-time feedback', 'dynamic UI',
    'session management'
]


# categorize
def categorize_difficulty(paragraph, keywords):
    """Categorize difficulty based on keyword presence and paragraph length."""
    length = len(paragraph)
    if any(keyword in paragraph.lower() for keyword in keywords['high_level']):
        return 'High Level'
    elif any(keyword in paragraph.lower() for keyword in keywords['medium_level']):
        return 'Medium Level'
    elif any(keyword in paragraph.lower() for keyword in keywords['low_level']):
        return 'Low Level'
    else:
        return 'Uncategorized'

# Function to extract from  document
def extract_paragraphs_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    paragraphs = []
    for page in doc:
        text = page.get_text("text")
        page_paragraphs = text.split('\n')
        cleaned_paragraphs = []
        paragraph = ""
        for line in page_paragraphs:
            if line.strip():
                paragraph += line.strip() + " "
            else:
                if paragraph:
                    cleaned_paragraphs.append(paragraph.strip())
                    paragraph = ""
        if paragraph:
            cleaned_paragraphs.append(paragraph.strip())
        paragraphs.extend(cleaned_paragraphs)
    return paragraphs

def extract_paragraphs_from_word(word_path):
    doc = Document(word_path)
    return [para.text for para in doc.paragraphs if para.text.strip()]

def extract_paragraphs_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    paragraphs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                paragraphs.append(shape.text.strip())
    return paragraphs

# Function to generate question
def generate_question_and_options(paragraph, keywords):
    doc = nlp(paragraph)

    # Extract meaningful data
    nouns = [token.text for token in doc if token.pos_ in ['NOUN', 'PROPN']]
    question = "What is the main idea of this paragraph?"
    correct_answer = paragraph[:100] + "..." if len(paragraph) > 100 else paragraph

    if nouns:
        noun_phrase = " ".join(nouns[:2]) if len(nouns) > 1 else nouns[0]
        question = f"What is discussed about {noun_phrase}?"

    sentences = [sent.text for sent in doc.sents]
    incorrect_answers = random.sample(sentences[1:], min(3, len(sentences) - 1)) if len(sentences) > 1 else []

    while len(incorrect_answers) < 3:
        incorrect_answers.append("Random incorrect statement.")

    options = [correct_answer] + incorrect_answers
    random.shuffle(options)

    # Categorize
    difficulty = categorize_difficulty(paragraph, keywords)

    return {"question": question, "options": options, "correct_answer": correct_answer, "difficulty": difficulty }

# upload_file
@app.route('/upload', methods=['POST'])
def upload_file():
    global questions_store
    if 'file' not in request.files:
        return jsonify({"error": "No file part"}), 400

    file = request.files['file']
    if file.filename == '':
        return jsonify({"error": "No selected file"}), 400

    file_path = os.path.join(UPLOAD_FOLDER, file.filename)
    file.save(file_path)

    file_ext = file.filename.split('.')[-1].lower()
    if file_ext == 'pdf':
        paragraphs = extract_paragraphs_from_pdf(file_path)
    elif file_ext == 'docx':
        paragraphs = extract_paragraphs_from_word(file_path)
    elif file_ext == 'txt':
        with open(file_path, 'r') as f:
            paragraphs = [line.strip() for line in f.readlines() if line.strip()]
    elif file_ext in ['ppt', 'pptx']:
        paragraphs = extract_paragraphs_from_ppt(file_path)
    else:
        return jsonify({"error": "Unsupported file type"}), 400

    generated_questions = [generate_question_and_options(p, {
        'high_level': high_level_keywords,
        'medium_level': medium_level_keywords,
        'low_level': low_level_keywords
    }) for p in paragraphs[:5]]

    with open("results_summary.txt", "r") as file:
        level_from_file = file.read().strip()

    questions_store = [q for q in generated_questions if q['difficulty'] == level_from_file]

    return jsonify({"message": "Questions generated successfully"})

@app.route('/get_questions', methods=['GET'])
def get_questions():
    with open("results_summary.txt", "r") as file:
        level_from_file = file.read().strip()
        print(level_from_file)

    return jsonify({"questions": questions_store , "level": level_from_file })


@app.route('/submit_answers', methods=['POST'])
def submit_answers():
    """Evaluate answers and return results."""
    data = request.json
    user_answers = data.get('answers', {})

    correct_count = 0
    incorrect_count = 0
    results = []

    # Iterate 
    for idx, question in enumerate(questions_store):
        correct_answer = question["correct_answer"]
        user_answer = user_answers.get(str(idx), "")
        result = "Correct" if user_answer == correct_answer else "Incorrect"

        # Count correct and incorrect answers
        if result == "Correct":
            correct_count += 1
        else:
            incorrect_count += 1

        results.append({
            "question": question["question"],
            "user_answer": user_answer,
            "correct_answer": correct_answer,
            "result": result,
            "difficulty": question["difficulty"]
        })

    #  percentage
    total_answers = len(questions_store)
    correct_percentage = (correct_count / total_answers) * 100 if total_answers > 0 else 0


    if correct_percentage > 75:
        level = "High Level"
    elif 45 <= correct_percentage <= 75:
        level = "Medium Level"
    else:
        level = "Low Level"

    # Save 
    with open("results_summary.txt", "w") as file:
        file.write(f"{level}\n")

    # Read the level 
    with open("results_summary.txt", "r") as file:
        level_from_file = file.read().strip()


    return jsonify({
        "results": results,
        "correct_count": correct_count,
        "incorrect_count": incorrect_count,
        "correct_percentage": correct_percentage,
        "level": level_from_file
    })


@app.route('/get_level', methods=['GET'])
def get_level():
    try:
        with open("results_summary.txt", "r") as file:
            level = file.read().strip()
        return jsonify({"level": level})
    except Exception as e:
        return jsonify({"error": "Unable to read level from file"}), 500


if __name__ == '__main__':
    app.run(debug=True)
