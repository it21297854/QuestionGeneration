import fitz
from docx import Document
from pptx import Presentation
import re

def extract_paragraphs_from_pdf(pdf_path):
    doc = fitz.open(pdf_path)
    paragraphs = []
    for page in doc:
        text = page.get_text("text")
        page_paragraphs = text.split('\n')
        cleaned_paragraphs = []
        paragraph = ""
        for line in page_paragraphs:
            if line.strip():
                paragraph += line.strip() + " "
            else:
                if paragraph:
                    cleaned_paragraphs.append(paragraph.strip())
                    paragraph = ""
        if paragraph:
            cleaned_paragraphs.append(paragraph.strip())
        paragraphs.extend(cleaned_paragraphs)
    return paragraphs

def extract_paragraphs_from_word(word_path):
    doc = Document(word_path)
    return [para.text for para in doc.paragraphs if para.text.strip()]

def extract_paragraphs_from_ppt(ppt_path):
    prs = Presentation(ppt_path)
    paragraphs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                paragraphs.append(shape.text.strip())
    return paragraphs

def clean_text(paragraph):
    # Remove page numbers, chapter numbers, headers, and unwanted leading numbers
    paragraph = re.sub(r'\b(Page \d+|Chapter \d+|Header|^\d+\.?)\b', '', paragraph)  # Remove page numbers, chapter numbers, headers, leading numbers
    paragraph = re.sub(r'\s+', ' ', paragraph).strip()  # Remove extra spaces and newlines
    paragraph = re.sub(r'\d+', '', paragraph)  # Remove any standalone digits (like '1', '2', etc.)
    return paragraph

